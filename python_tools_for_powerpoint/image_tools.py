import os
from pptx.util import Inches
from stat import S_ISREG, ST_CTIME, ST_MODE
from python_tools_for_powerpoint import pptx_tools


def find_images(folder_path, extensions=('.png', '.jpg', '.tif')):
    """
    This function searches a folder for files that have one of the given file extensions.

    :param folder_path: a string path to the images folder
    :param extensions: a list of image extensions strings like ['.png', '.jpg']
    :return: a list of image file names strings
    """
    image_full_paths = []
    for filename in os.listdir(folder_path):
        basename, extension = os.path.splitext(filename)
        if extension.lower() in extensions:
            image_full_paths.append(os.path.join(folder_path, filename))
    return image_full_paths


def sort_files_by_creation_date(path_list):
    data = ((os.stat(path), path) for path in path_list)
    data = ((stat[ST_CTIME], path) for stat, path in data if S_ISREG(stat[ST_MODE]))
    data = sorted(data)
    return [path for stat, path in data]


def add_full_slide_image(prs, slide_obj, img_path):
    slide_obj.shapes.add_picture(img_path, left=Inches(0), top=Inches(0), width=prs.slide_width)


def add_image_slide(prs, img_path):
    """
    This function creates a new blank slide and adds a image to it. The image is rescaled to
    fit the width of the slide, while retaining its aspect ratio.

    :param prs: the Presentation object
    :param img_path: a string containing the path to the image file
    :return: a Presentation object
    """
    # Add a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Add the image, resizing it to fit slide width, while maintaining aspect ratio
    new_slide.shapes.add_picture(img_path, left=Inches(0), top=Inches(0), width=prs.slide_width)
    return prs


def add_multiple_image_slides_from_folder(prs, img_folder, extensions=('.png', '.jpg', '.tif')):
    for img_path in find_images(img_folder, extensions):
        add_image_slide(prs, img_path)
        print("Added image slide containing image: '{0}'".format(img_path))


def add_multiple_image_slide_from_list_with_notes(prs, img_path_list, notes_generator):
    for img_path in img_path_list:
        slide_obj = pptx_tools.insert_blank_slide(prs)
        add_full_slide_image(prs, slide_obj, img_path)
        slide_notes = notes_generator(img_path)

        if slide_notes is not None:
            # Assume notes are formatted object
            try:
                note_run_text = slide_notes[0][0]['text']
                pptx_tools.add_formatted_notes(slide_notes, slide_obj)
            # Assume notes are string
            except TypeError:
                pptx_tools.add_notes_to_slide(slide_notes, slide_obj)

    return prs


def add_multiple_image_slide_from_list(prs, img_path_list):
    for img_path in img_path_list:
        add_image_slide(prs, img_path)


def add_multiple_image_slides_from_folder_by_date(prs, img_folder):
    data = (os.path.join(img_folder, fn) for fn in os.listdir(img_folder))
    data = ((os.stat(path), path) for path in data)
    data = ((stat[ST_CTIME], path) for stat, path in data if S_ISREG(stat[ST_MODE]))
    for cdate, img_path in sorted(data):
        add_image_slide(prs, img_path)
