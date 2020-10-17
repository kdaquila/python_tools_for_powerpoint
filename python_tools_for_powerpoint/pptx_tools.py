import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def insert_blank_slide(prs):
    # Add a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_slide_layout)


def load_pptx_file(pptx_file_path):
    # Load existing presentation
    try:
        prs = pptx.Presentation(pptx_file_path)
        print("\nFound existing file, will copy and edit")
        return prs
    except pptx.exc.PackageNotFoundError:
        print("No existing file found")
        return None


def add_notes_to_slide(notes_text, slide_obj):
    slide_obj.notes_slide.notes_text_frame.text = notes_text


def insert_slide_notes(notes_text, prs, slide_num=0):
    prs.slides[slide_num].notes_slide.notes_text_frame.text = notes_text
    return prs


def add_formatted_notes(notes_obj, slide_obj):
    text_frame = slide_obj.notes_slide.notes_text_frame
    is_first_paragraph = True
    for note_paragraph in notes_obj:
        if is_first_paragraph:
            paragraph = text_frame.paragraphs[0]
            is_first_paragraph = False
        else:
            paragraph = text_frame.add_paragraph()
        for note_run in note_paragraph:
            run = paragraph.add_run()
            run.text = note_run['text']
            if note_run.get('bold', False):
                run.font.bold = True


def find_first_table(prs, slide_num):
    slide = prs.slides[slide_num]
    for shape in slide.shapes:
        print(shape.shape_type)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape
    raise Exception("The slide does not contain a table")


def apply_conditional_formatting_to_columns(prs, slide_num, shape_id, col_rules):
    # Find the powerpoint table
    slide = prs.slides[slide_num]
    table = slide.shapes[shape_id]

    # Get the table dimensions
    num_rows = len(table.table.rows)
    num_cols = len(table.table.columns)

    # Visit all columns that have a rule defined
    for rule in col_rules:
        col_ind = rule['column']

        # Visit all rows
        for row_ind in range(num_rows):

            # Read the cell value
            cell_value_str = table.table.cell(row_ind, col_ind).text

            # Compute the conditional color
            conditional_color = rule['func'](cell_value_str)

            # Apply the color to all paragraphs and runs
            if conditional_color is not None:
                for para in table.table.cell(row_ind, col_ind).text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor.from_string(conditional_color)
