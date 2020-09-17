import os
from pptx.util import Inches


def find_images(folder_path, extensions=('.png', '.jpg', '.tif')):
    """
    This function searches a folder for files that have one of the given file extensions.

    :param folder_path: a string path to the images folder
    :param extensions: a list of image extensions strings like ['.png', '.jpg']
    :return: a list of image file names strings
    """
    image_full_paths = []
    for filename in os.listdir(folder_path):
        basename, extension = os.path.splitext(filename)
        if extension.lower() in extensions:
            image_full_paths.append(os.path.join(folder_path, filename))
    return image_full_paths


def add_image_slide(prs, img_path):
    """
    This function creates a new blank slide and adds a image to it. The image is rescaled to
    fit the width of the slide, while retaining its aspect ratio.

    :param prs: the Presentation object
    :param img_path: a string containing the path to the image file
    :return: a Presentation object
    """
    # Add a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Add the image, resizing it to fit slide width, while maintaining aspect ratio
    new_slide.shapes.add_picture(img_path, left=Inches(0), top=Inches(0), width=prs.slide_width)
    return prs


def add_multiple_image_slides(prs, img_folder, extensions=('.png', '.jpg', '.tif')):
    for img_path in find_images(img_folder, extensions):
        add_image_slide(prs, img_path)
        print("Added image slide containing image: '{0}'".format(img_path))
