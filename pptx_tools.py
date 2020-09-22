import pptx


def load_pptx_file(pptx_file_path):
    # Load existing presentation
    try:
        prs = pptx.Presentation(pptx_file_path)
        print("Found existing file, will copy and edit")
        return prs
    except pptx.exc.PackageNotFoundError:
        print("No existing file found")
        return None


def insert_slide_notes(prs, slide_num, notes_text):
    prs.slides[slide_num].notes_slide.notes_text_frame.text = notes_text
    return prs


