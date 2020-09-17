from pptx.util import Inches
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.util import Pt


def add_textbox_horiz_align_center(prs, slide_num, text_string, top_inch=4.4,
                                   font_size_pt=24, font_name='Calabri'):

    # Insert the text box
    text_shape = prs.slides[slide_num].shapes.add_textbox(left=Inches(0),
                                                          top=Inches(top_inch),
                                                          width=Inches(prs.slide_width.inches),
                                                          height=Inches(font_size_pt / 72.0 * 1.2))

    # Set the text string
    text_frame = text_shape.text_frame
    text_paragraph = text_frame.paragraphs[0]
    text_run = text_paragraph.add_run()
    text_run.text = text_string

    # Set center alignment
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    text_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Set the font
    font = text_run.font
    font.name = font_name
    font.size = Pt(font_size_pt)
    font.bold = False
    font.italic = False
    font.color.theme_color = MSO_THEME_COLOR_INDEX.TEXT_1

    return prs
