import os
import pptx
import re
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import openpyxl
from openpyxl.styles.colors import COLOR_INDEX



def open_excel_file(excel_file_path):
    return openpyxl.load_workbook(filename=excel_file_path)


def load_office_theme_colors():
    return [
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': 0.0, 'hex': 'FFFFFF'},
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': -0.05, 'hex': 'F2F2F2'},
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': -0.15, 'hex': 'D9D9D9'},
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': -0.25, 'hex': 'BFBFBF'},
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': -0.35, 'hex': 'A6A6A6'},
        {'theme_name': 'BACKGROUND_1', 'theme_id': 0, 'tint': -0.5, 'hex': '808080'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.0, 'hex': '000000'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.5, 'hex': '808080'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.35, 'hex': '595959'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.25, 'hex': '404040'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.15, 'hex': '262626'},
        {'theme_name': 'TEXT_1', 'theme_id': 1, 'tint': 0.05, 'hex': '0D0D0D'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': 0.0, 'hex': 'E7E6E6'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': -0.1, 'hex': 'D0CECE'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': -0.25, 'hex': 'AEAAAA'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': -0.5, 'hex': '757171'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': -0.75, 'hex': '3A3838'},
        {'theme_name': 'BACKGROUND_2', 'theme_id': 2, 'tint': -0.9, 'hex': '161616'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': 0.0, 'hex': '44546A'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': 0.8, 'hex': 'D6DCE4'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': 0.6, 'hex': 'ACB9CA'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': 0.4, 'hex': '8497B0'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': -0.25, 'hex': '333F4F'},
        {'theme_name': 'TEXT_2', 'theme_id': 3, 'tint': -0.5, 'hex': '222B35'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': 0.0, 'hex': '4472C4'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': 0.8, 'hex': 'D9E1F2'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': 0.6, 'hex': 'B4C6E7'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': 0.4, 'hex': '8EA9DB'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': -0.25, 'hex': '305496'},
        {'theme_name': 'ACCENT_1', 'theme_id': 4, 'tint': -0.5, 'hex': '203764'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': 0.0, 'hex': 'ED7D31'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': 0.8, 'hex': 'FCE4D6'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': 0.6, 'hex': 'F8CBAD'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': 0.4, 'hex': 'F4B084'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': -0.25, 'hex': 'C65911'},
        {'theme_name': 'ACCENT_2', 'theme_id': 5, 'tint': -0.5, 'hex': '833C0C'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': 0.0, 'hex': 'A5A5A5'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': 0.8, 'hex': 'EDEDED'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': 0.6, 'hex': 'DBDBDB'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': 0.4, 'hex': 'C9C9C9'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': -0.25, 'hex': '7B7B7B'},
        {'theme_name': 'ACCENT_3', 'theme_id': 6, 'tint': -0.5, 'hex': '525252'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': 0.0, 'hex': 'FFC000'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': 0.8, 'hex': 'FFF2CC'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': 0.6, 'hex': 'FFE699'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': 0.4, 'hex': 'FFD966'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': -0.25, 'hex': 'BF8F00'},
        {'theme_name': 'ACCENT_4', 'theme_id': 7, 'tint': -0.5, 'hex': '806000'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': 0.0, 'hex': '5B9BD5'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': 0.8, 'hex': 'DDEBF7'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': 0.6, 'hex': 'BDD7EE'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': 0.4, 'hex': '9BC2E6'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': -0.25, 'hex': '2F75B5'},
        {'theme_name': 'ACCENT_5', 'theme_id': 8, 'tint': -0.5, 'hex': '1F4E78'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': 0.0, 'hex': '70AD47'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': 0.8, 'hex': 'E2EFDA'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': 0.6, 'hex': 'C6E0B4'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': 0.4, 'hex': 'A9D08E'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': -0.25, 'hex': '548235'},
        {'theme_name': 'ACCENT_6', 'theme_id': 9, 'tint': -0.5, 'hex': '375623'}]


def find_theme_color_hex(theme_id, tint, theme_colors):
    for color in theme_colors:
        if color['theme_id'] == theme_id and color['tint'] == round(tint, 2):
            return color['hex']
    return None


def is_indexable(obj):
    try:
        len(obj)
        return True
    except TypeError:
        return False


def is_iterable(obj):
    try:
        iter(obj)
        return True
    except TypeError:
        return False


def is_sheet_cell_ref(obj):
    if is_indexable(obj) and is_iterable(obj) and obj[0] == '=' and "!" in obj:
        return True
    else:
        return False


def get_rbg_from_color_obj(color_obj, theme_colors):
    if color_obj is None:
        return None
    color_type = color_obj.type
    if color_type == 'theme':
        theme_val = color_obj.theme
        tint_val = color_obj.tint
        return find_theme_color_hex(theme_val, tint_val, theme_colors)
    elif color_type == 'indexed':
        return openpyxl.styles.colors.COLOR_INDEX[color_obj.value - 1]
    elif color_type == 'rgb':
        argb_str = color_obj.value
        rgb_str = argb_str[2:]
        return rgb_str


def get_cell_value(wb, sheet, cell):
    cell_raw_value = wb[sheet][cell].value

    if cell_raw_value is None:
        return "None"

    if is_sheet_cell_ref(cell_raw_value):
        (sheet, cell) = cell_raw_value.strip("=").split("!")
        sheet = sheet.strip("'")

    return str(wb[sheet][cell].value)


def get_cell_alignment(wb, sheet, cell):
    return wb[sheet][cell].alignment.horizontal


def get_cell_font_color(wb, sheet, cell, theme_colors):
    return get_rbg_from_color_obj(wb[sheet][cell].font.color, theme_colors)


def get_cell_font_size(wb, sheet, cell):
    return wb[sheet][cell].font.sz


def get_cell_fill_color(wb, sheet, cell, theme_colors):
    return get_rbg_from_color_obj(wb[sheet][cell].fill.fgColor, theme_colors)


def get_conditional_color(value_str, col):

    if col == 1:
        if re.fullmatch('[0-9]+ Days', value_str) is None:
            return None
        num_days = int(value_str.split(" ")[0])
        if num_days >= 7:
            return 'FF0000'
        elif num_days >= 5:
            return '0070C0'

    elif col == 2:
        if re.fullmatch('[0-9]+ Days', value_str) is None:
            return None
        num_days = int(value_str.split(" ")[0])
        if num_days >= 3:
            return 'FF0000'
        elif num_days >= 2:
            return '0070C0'

    elif col in (3, 4):
        if re.fullmatch('[0-9]+.[0-9]+%', value_str) is None:
            return None
        float_value = float(value_str.split("%")[0])
        if float_value >= 8.0:
            return 'FF0000'
        elif float_value >= 7.0:
            return '0070C0'

    elif col == 5:
        if re.fullmatch('[0-9]+ Days', value_str) is None:
            return None
        num_days = int(value_str.split(" ")[0])
        if num_days >= 7:
            return 'FF0000'
        elif num_days >= 5:
            return '0070C0'

    return None


def insert_excel_range(prs, excel_file_path, sheet=0, col_start='A', col_stop='B', row_start=1, row_stop=2,
                       slide_num=0, top_inch=1, left_inch=1, width_inch=3, height_inch=3, font_size_factor=1.0):
    # Open the Excel File
    wb = open_excel_file(excel_file_path)

    # Load the office theme colors
    office_theme_colors = load_office_theme_colors()

    # Create the empty powerpoint table
    slide = prs.slides[slide_num]
    num_row = row_stop - row_start + 1
    num_col = ord(col_stop) - ord(col_start) + 1
    table = slide.shapes.add_table(rows=num_row, cols=num_col, left=Inches(left_inch), top=Inches(top_inch),
                                   width=Inches(width_inch), height=Inches(height_inch))

    # Turn off special first-row formatting
    table.table.first_row = False

    # Set minimum row heights
    for row in table.table.rows:
        row.height = 0

    # Merge ranges as needed
    for merge_range in wb[sheet].merged_cells.ranges:

        merge_range_bounds = merge_range.bounds
        top_left_col = merge_range_bounds[0] - 1 - (ord(col_start) - ord('A'))
        top_left_row = merge_range_bounds[1] - row_start
        bottom_right_col = merge_range_bounds[2] - 1 - (ord(col_start) - ord('A'))
        bottom_right_row = merge_range_bounds[3] - row_start

        if (top_left_col >= 0 and
                top_left_row >= 0 and
                bottom_right_col <= (ord(col_stop) - ord(col_start)) and
                bottom_right_row <= (row_stop - row_start)):

            origin_cell = table.table.cell(top_left_row, top_left_col)
            other_cell = table.table.cell(bottom_right_row, bottom_right_col)
            origin_cell.merge(other_cell)

    # Fill the powerpoint table cell-by-cell
    for row in range(row_start - row_start, row_stop - row_start + 1):
        for col in range(0, ord(col_stop) - ord(col_start) + 1):
            excel_col_letter = chr((col + ord(col_start) - ord('A')) + ord('A'))
            excel_row_num = row + row_start
            cell_ref = str(excel_col_letter) + str(excel_row_num)

            # Get the excel cell value
            cell_value = get_cell_value(wb, sheet, cell_ref)

            # Set the cell value
            table.table.cell(row, col).text = cell_value

            # Get the excel cell font color
            cell_font_color = get_cell_font_color(wb, sheet, cell_ref, office_theme_colors)

            # Get the excel cell's conditional color
            cell_cond_font_color = get_conditional_color(cell_value, col)
            if cell_cond_font_color is not None:
                cell_font_color = cell_cond_font_color

            # Transfer font color and set font weight to normal
            for para in table.table.cell(row, col).text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor.from_string(cell_font_color)
                    # run.font.bold = False

            # Get the excel cell fill color
            cell_fill_color = get_cell_fill_color(wb, sheet, cell_ref, office_theme_colors)

            # Transfer fill color
            table.table.cell(row, col).fill.solid()
            table.table.cell(row, col).fill.fore_color.rgb = RGBColor.from_string(cell_fill_color)

            # Transfer the font size
            cell_font_size = get_cell_font_size(wb, sheet, cell_ref)
            for para in table.table.cell(row, col).text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(int(round(cell_font_size * font_size_factor)))

            # Get cell alignment
            cell_alignment = get_cell_alignment(wb, sheet, cell_ref)

            # Set cell alignment
            for para in table.table.cell(row, col).text_frame.paragraphs:
                if cell_alignment is None:
                    pass
                elif cell_alignment.upper() == 'LEFT':
                    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                elif cell_alignment.upper() == 'RIGHT':
                    para.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                elif cell_alignment.upper() == 'CENTER':
                    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    return prs
