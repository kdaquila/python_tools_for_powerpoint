import xlwings as xw
import PIL.ImageGrab
import os
from pptx.util import Inches


def delete_temp_image(folder_path):
    for filename in os.listdir(folder_path):
        if filename == 'temp.png':
            os.remove(os.path.join(folder_path, filename))


def close_excel():
    app = xw.apps.active
    app.quit()


def excel_to_img(excel_file_path, temp_img_path, sheet_num=0, cell_range='A1:B1'):
    book = xw.Book(excel_file_path)
    sheet = book.sheets[sheet_num]
    sheet.range(cell_range).api.CopyPicture(Appearance=1, Format=2)
    sheet.api.Paste()
    pic = sheet.pictures[0]
    pic.api.Copy()
    img = PIL.ImageGrab.grabclipboard()
    img.save(temp_img_path)


def insert_excel_as_img(prs, excel_file_path, sheet=0, cell_range='A1:B2', slide_num=0, top_inch=1, left_inch=1):

    # Find the excel file's directory as use it for temp img
    (temp_img_dir, excel_file_name) = os.path.split(excel_file_path)

    # Build path to temp.png
    temp_img_path = os.path.join(temp_img_dir, 'temp.png')

    # Generate the temp.png image from excel range
    excel_to_img(excel_file_path, temp_img_path, sheet, cell_range)
    close_excel()

    # Insert the picture
    prs.slides[slide_num].shapes.add_picture(temp_img_path, left=Inches(left_inch), top=Inches(top_inch),
                                             width=Inches(prs.slide_width.inches - 2 * left_inch))

    # Delete the temp image
    os.remove(temp_img_path)

    return prs


